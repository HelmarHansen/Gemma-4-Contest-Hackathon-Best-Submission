import json
import re
import base64
import io
import dataclasses
import requests
from dataclasses import dataclass, field, asdict
from fastapi import FastAPI, HTTPException
from fastapi.responses import StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

app = FastAPI()

OLLAMA_URL = "http://localhost:11434/api/chat"
MODEL      = "gemma4"

# ── Ollama parameter presets ─────────────────────────────────────
# num_predict is sized so that the longer worked examples and the
# 7-CHECK-validated blueprints from setup_prompt.txt cannot be truncated.
# Lifted from 512/6000/512/1500 → 2048/12288/1536/2048.
_OPTS_RESEARCH    = {"temperature": 0.15, "top_p": 0.85, "top_k": 30,
                     "repeat_penalty": 1.05, "num_predict": 2048}
_OPTS_TASK_DESIGN = {"temperature": 0.25, "top_p": 0.85, "top_k": 40,
                     "repeat_penalty": 1.10, "num_predict": 8192}
_OPTS_BLUEPRINT   = {"temperature": 0.25, "top_p": 0.85, "top_k": 40,
                     "repeat_penalty": 1.10, "num_predict": 12288}
_OPTS_VALIDATE    = {"temperature": 0.10, "top_p": 0.80, "top_k": 20,
                     "repeat_penalty": 1.00, "num_predict": 1536}
_OPTS_CHAT        = {"temperature": 0.30, "top_p": 0.88, "top_k": 45,
                     "repeat_penalty": 1.20, "num_predict": 2048}
_OPTS_INTERROGATE = {"temperature": 0.45, "top_p": 0.90, "top_k": 50,
                     "repeat_penalty": 1.15, "num_predict": 768}

# ── Mood mapping ─────────────────────────────────────────────────

_MOOD_BY_MODE: dict[str, list[str]] = {
    "Cold Case":  ["calm",    "calm",    "neutral", "tense",  "calm"   ],
    "Murder":     ["neutral", "warm",    "tense",   "danger", "danger" ],
    "Conspiracy": ["neutral", "neutral", "tense",   "danger", "tense"  ],
    "Heist":      ["warm",    "warm",    "tense",   "tense",  "neutral"],
}
_MOOD_DEFAULT = ["neutral", "warm", "tense", "tense", "danger"]

def _mood_for(mode: str, phase_idx: int) -> str:
    phases = _MOOD_BY_MODE.get(mode, _MOOD_DEFAULT)
    return phases[min(phase_idx, len(phases) - 1)]

# ── Blueprint dataclasses ────────────────────────────────────────

@dataclass
class Character:
    id:               str = ""
    name:             str = ""
    role:             str = ""
    personality:      str = ""
    voice:            str = ""
    knowledge:        str = ""
    secrets:          str = ""
    ignorant_of:      str = ""
    appears_in_moves: list[str] = field(default_factory=list)

@dataclass
class Teacher:
    name:                 str = ""
    role:                 str = ""
    voice:                str = ""
    active_traits:        list[str] = field(default_factory=list)
    expects_from_student: str = ""
    will_not_do:          str = ""
    signature_phrases:    list[str] = field(default_factory=list)

@dataclass
class Session:
    mode:            str   = ""
    language:        str   = "English"
    immersive:       bool  = True
    difficulty:      float = 0.5
    pacing:          str   = "steady"
    estimated_turns: int   = 25
    hint_policy:     str   = "only if asked"

@dataclass
class Topic:
    raw:               str = ""
    core_subject:      str = ""
    learning_goals:    list[str] = field(default_factory=list)
    likely_weak_spots: list[str] = field(default_factory=list)
    out_of_scope:      str = ""
    material_summary:  str = ""

@dataclass
class Phase:
    id:               str = ""
    name:             str = ""
    goal:             str = ""
    turn_budget:      int = 5
    teacher_strategy: str = ""
    exit_condition:   str = ""
    fallback:         str = ""

@dataclass
class Move:
    id:                        str = ""
    phase:                     str = ""
    type:                      str = "scene"
    trigger:                   str = ""
    content:                   str = ""
    expected_student_response: str = ""
    if_correct:                str = ""
    if_incorrect:              str = ""
    hint:                      str | None = None
    source_ref:                str | None = None 
    is_error_trap:             bool       = False
    error_correction_key:      str | None = None

@dataclass
class Contingencies:
    student_goes_off_topic:        str = ""
    student_asks_for_answer:       str = ""
    student_is_silent:             str = ""
    student_is_consistently_wrong: str = ""
    student_finishes_early:        str = ""

@dataclass
class Opening:
    first_line:      str = ""
    opening_move_id: str = ""

@dataclass
class Closing:
    success_line: str = ""
    partial_line: str = ""
    fail_line:    str = ""

@dataclass
class ExecutionRules:
    how_to_pick_next_move:       str = ""
    when_to_deviate_from_plan:   str = ""
    response_length_guideline:   str = ""
    never_do:                    list[str] = field(default_factory=list)

@dataclass
class SessionBlueprint:
    session_id:      str
    teacher:         Teacher
    session:         Session
    topic:           Topic
    phases:          list[Phase]
    moves:           list[Move]
    characters:      list[Character]
    contingencies:   Contingencies
    opening:         Opening
    closing:         Closing
    execution_rules: ExecutionRules
    _move_index: dict[str, Move]      = field(default_factory=dict, repr=False)
    _char_index: dict[str, Character] = field(default_factory=dict, repr=False)

    def __post_init__(self):
        self._move_index = {m.id: m for m in self.moves}
        self._char_index = {c.id: c for c in self.characters}

    def get_move(self, move_id: str) -> Move | None:
        return self._move_index.get(move_id)

    def get_character(self, char_id: str) -> Character | None:
        return self._char_index.get(char_id)

    def get_opening_move(self) -> Move | None:
        return self.get_move(self.opening.opening_move_id)

    def to_dict(self) -> dict:
        d = asdict(self)
        d.pop("_move_index", None)
        d.pop("_char_index", None)
        return d

    @classmethod
    def _from_dict(cls, data: dict) -> "SessionBlueprint":
        def fit(cls_, d):
            # Strip unknown keys so new blueprint fields don't break old dataclasses
            if not isinstance(d, dict):
                return cls_()
            keys = {f.name for f in dataclasses.fields(cls_)}
            return cls_(**{k: v for k, v in d.items() if k in keys})

        def get(key, default=None):
            return data.get(key) or default

        return cls(
            session_id=get("session_id", "session"),
            teacher=fit(Teacher, get("teacher", {})),
            session=fit(Session, get("session", {})),
            topic=fit(Topic, get("topic", {})),
            phases=[fit(Phase, p) for p in get("arc", {}).get("phases", [])],
            moves=[fit(Move, m) for m in get("planned_moves", [])],
            characters=[fit(Character, c) for c in get("characters", [])],
            contingencies=fit(Contingencies, get("contingencies", {})),
            opening=fit(Opening, get("opening", {})),
            closing=fit(Closing, get("closing", {})),
            execution_rules=fit(ExecutionRules, get("execution_rules", {})),
        )

    @classmethod
    def from_llm_response(cls, raw: str) -> "SessionBlueprint":
        return cls._from_dict(json.loads(_extract_json_object(raw)))

# ── Request schema ───────────────────────────────────────────────

class TeacherInput(BaseModel):
    name:        str
    role:        str
    personality: str
    traits:      list[str]

class LessonInput(BaseModel):
    topic:       str
    mode:        str
    language:    str
    length:      str
    difficulty:  float
    school_type: str = ""
    grade:       str = ""

class DocumentInput(BaseModel):
    name: str
    data: str  # base64-encoded file bytes

class WorkRequest(BaseModel):
    teacher:   TeacherInput
    lesson:    LessonInput
    material:  str = ""
    images:    list[str] = []           
    documents: list[DocumentInput] = []

class ChatMessage(BaseModel):
    role:    str
    content: str

class ChatRequest(BaseModel):
    session_id: str
    message:    str
    history:    list[ChatMessage] = []

class InterrogationRequest(BaseModel):
    session_id:   str
    character_id: str
    message:      str
    history:      list[ChatMessage] = []

# ── Ollama helpers ───────────────────────────────────────────────

def ask_ollama(system_prompt: str, user_message: str,
               options: dict | None = None,
               images: list[str] | None = None,
               force_json: bool = False) -> str:
    user_msg: dict = {"role": "user", "content": user_message}
    if images:
        user_msg["images"] = images
    payload: dict = {
        "model": MODEL,
        "messages": [
            {"role": "system", "content": system_prompt},
            user_msg,
        ],
        "stream": False,
    }
    if options:
        payload["options"] = options
    if force_json:
        # Ollama's grammar-constrained JSON mode — generation is bounded so the
        # response is guaranteed to be syntactically valid JSON.
        payload["format"] = "json"
    response = requests.post(OLLAMA_URL, json=payload, timeout=600)
    response.raise_for_status()
    return response.json()["message"]["content"]

def ask_ollama_stream(system_prompt: str, user_message: str,
                      options: dict | None = None):
    payload: dict = {
        "model": MODEL,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user",   "content": user_message},
        ],
        "stream": True,
    }
    if options:
        payload["options"] = options
    with requests.post(OLLAMA_URL, json=payload, stream=True, timeout=300) as response:
        response.raise_for_status()
        for line in response.iter_lines():
            if line:
                data = json.loads(line)
                content = data.get("message", {}).get("content", "")
                if content:
                    yield content
                if data.get("done"):
                    break

def _load_prompt(filename: str) -> str:
    with open(filename, "r", encoding="utf-8") as f:
        return f.read()

# ── Document extraction ──────────────────────────────────────────

_MAX_PDF_PAGES   = 8
_PDF_RENDER_DPI  = 110
_MAX_DOC_CHARS   = 6000

def _extract_pdf(raw: bytes) -> tuple[str, list[str]]:
    """Render PDF pages to PNG (base64) + extract embedded text. Returns (text, images_b64)."""
    import fitz  # pymupdf
    text_chunks: list[str] = []
    images_b64: list[str]  = []
    with fitz.open(stream=raw, filetype="pdf") as doc:
        n_pages = min(doc.page_count, _MAX_PDF_PAGES)
        for i in range(n_pages):
            page = doc.load_page(i)
            t = str(page.get_text("text") or "")
            if t.strip():
                text_chunks.append(t.strip())
            pix = page.get_pixmap(dpi=_PDF_RENDER_DPI)
            images_b64.append(base64.b64encode(pix.tobytes("png")).decode("ascii"))
    return ("\n\n".join(text_chunks)[:_MAX_DOC_CHARS], images_b64)

def _extract_docx(raw: bytes) -> str:
    import docx
    d = docx.Document(io.BytesIO(raw))
    lines = [p.text for p in d.paragraphs if p.text.strip()]
    for tbl in d.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines)[:_MAX_DOC_CHARS]

def _extract_pptx(raw: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(raw))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                lines.append(text)
        if slide.has_notes_slide:
            notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
            note = str(getattr(notes_frame, "text", "") or "").strip()
            if note:
                lines.append(f"[Notes] {note}")
    return "\n".join(lines)[:_MAX_DOC_CHARS]

def _extract_xlsx(raw: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"--- Sheet: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                lines.append(" | ".join(cells))
            if len("\n".join(lines)) > _MAX_DOC_CHARS:
                break
    return "\n".join(lines)[:_MAX_DOC_CHARS]

def process_documents(docs: list["DocumentInput"]) -> tuple[str, list[str]]:
    """Dispatch by extension. Returns (extracted_text, extracted_images_b64)."""
    text_parts: list[str] = []
    images:     list[str] = []
    for doc in docs:
        try:
            raw = base64.b64decode(doc.data)
        except Exception as e:
            print(f"Skip {doc.name}: base64 decode failed ({e})")
            continue
        ext = doc.name.rsplit(".", 1)[-1].lower() if "." in doc.name else ""
        try:
            if ext == "pdf":
                txt, imgs = _extract_pdf(raw)
                if txt:
                    text_parts.append(f"--- {doc.name} ---\n{txt}")
                images.extend(imgs)
            elif ext == "docx":
                txt = _extract_docx(raw)
                if txt:
                    text_parts.append(f"--- {doc.name} ---\n{txt}")
            elif ext == "pptx":
                txt = _extract_pptx(raw)
                if txt:
                    text_parts.append(f"--- {doc.name} ---\n{txt}")
            elif ext == "xlsx":
                txt = _extract_xlsx(raw)
                if txt:
                    text_parts.append(f"--- {doc.name} ---\n{txt}")
            else:
                print(f"Skip {doc.name}: unsupported extension '{ext}'")
        except Exception as e:
            print(f"Extraction failed for {doc.name}: {e}")
    return ("\n\n".join(text_parts), images)

# ── Prompt building ──────────────────────────────────────────────

def _build_teacher_prompts(
    blueprint:        "SessionBlueprint",
    user_prompt:      str,
    current_move:     "Move | None",
    history:          list["ChatMessage"] | None = None,
    visited_move_ids: list[str] | None = None,
    accuracy_ratio:   float = 0.5,
) -> tuple[str, str]:
    """Returns (system_prompt, user_message) for the teacher agent."""
    system_template = _load_prompt("teacher_system_prompt.txt")
    system_prompt = system_template.format(
        teacher_name=blueprint.teacher.name,
        teacher_role=blueprint.teacher.role,
        teacher_voice=blueprint.teacher.voice,
        teacher_traits=json.dumps(blueprint.teacher.active_traits, ensure_ascii=False),
        teacher_phrases=json.dumps(blueprint.teacher.signature_phrases, ensure_ascii=False),
        will_not_do=blueprint.teacher.will_not_do,
    )

    blueprint_json = json.dumps(blueprint.to_dict(), ensure_ascii=False, indent=2)

    history_text = "(none)"
    if history:
        lines = []
        for msg in history:
            label = "DETECTIVE" if msg.role == "user" else "GAME MASTER"
            lines.append(f"{label}: {msg.content}")
        history_text = "\n".join(lines)

    current_move_json = json.dumps(
        dataclasses.asdict(current_move) if current_move else {},
        ensure_ascii=False, indent=2,
    )
    is_closing_move = current_move and current_move.type == "closing"

    visited_ids = visited_move_ids or []
    already_tested_lines = []
    for mid in visited_ids:
        m = blueprint.get_move(mid)
        if m and m.id != (current_move.id if current_move else None):
            already_tested_lines.append(f"- [{m.id}] {m.expected_student_response}")
    already_tested = "\n".join(already_tested_lines) if already_tested_lines else "(none yet)"

    n = len(visited_move_ids or [])
    if n >= 3:
        if accuracy_ratio >= 0.80:
            accuracy_note = (
                f"HIGH accuracy ({accuracy_ratio:.0%}, {n} moves tested). "
                "Student answers are consistently correct — increase demand: require sub-processes, "
                "edge cases, or precise fine distinctions. Do not simplify or over-scaffold."
            )
        elif accuracy_ratio <= 0.25:
            accuracy_note = (
                f"LOW accuracy ({accuracy_ratio:.0%}, {n} moves tested). "
                "Student is struggling — embed a stronger environmental hint in [NARRATION]. "
                "In [NPC:Name] dialogue, explicitly name both the student's wrong term and the "
                "correct term, with one clause stating how they differ."
            )
        else:
            accuracy_note = (
                f"BALANCED ({accuracy_ratio:.0%}, {n} moves tested). Maintain current complexity."
            )
    else:
        accuracy_note = "Too early to calibrate (fewer than 3 moves tested). Maintain baseline complexity from the case file."

    user_template = _load_prompt("teacher_user_prompt.txt")
    user_message = user_template.format(
        blueprint_json=blueprint_json,
        history=history_text,
        user_prompt=user_prompt,
        current_move_json=current_move_json,
        is_closing_move="YES — this is the final move. Use CLOSE if the student answers correctly."
                        if is_closing_move else "NO",
        already_tested=already_tested,
        accuracy_note=accuracy_note,
    )
    return system_prompt, user_message

def run_teacher_agent(
    blueprint:        "SessionBlueprint",
    user_prompt:      str,
    current_move:     "Move | None",
    history:          list["ChatMessage"] | None = None,
    visited_move_ids: list[str] | None = None,
    accuracy_ratio:   float = 0.5,
) -> tuple[str, str]:
    """Returns (clean_reply_without_STATE, state_signal)."""
    system_prompt, user_message = _build_teacher_prompts(
        blueprint, user_prompt, current_move, history, visited_move_ids, accuracy_ratio,
    )
    raw = ask_ollama(system_prompt, user_message, options=_OPTS_CHAT)

    state = "STAY"
    state_match = re.search(
        r"\[STATE\][:\s]*\n?\s*(ADVANCE|STAY|CLOSE)", raw, re.IGNORECASE
    )
    if state_match:
        state = state_match.group(1).upper()
        raw = raw[:state_match.start()].rstrip()

    return raw, state

# ── Validation system prompt ─────────────────────────────────────

_VALIDATION_SYSTEM = """\
You are a factual accuracy checker for educational content.
You receive numbered "expected student answers" from a detective learning scenario, each labelled with its subject.
For each entry, check whether the stated fact or mechanism is scientifically / historically / mathematically correct.
Output exactly one line per entry:
  OK: [number]
  FLAG: [number] — [one sentence explaining the factual error]
No other text. Be strict: flag anything that is wrong, oversimplified to the point of being misleading, or not verifiable.
"""

# ── Three-stage story builder ────────────────────────────────────

def _extract_json_object(raw: str) -> str:
    """Salvage a JSON object from messy LLM output.
    Handles: markdown code fences, smart quotes, raw backslashes (`\\frac`,
    Windows paths), raw newlines/tabs inside strings, and trailing commas.
    Trims to the first string-aware balanced `{...}` so trailing prose doesn't
    poison `json.loads`.
    """
    # Strip markdown code fences (```json ... ``` or just ```).
    s = re.sub(r"```(?:json|JSON)?\s*", "", raw)
    s = s.replace("```", "")
    # Normalize smart quotes to ASCII.
    s = (s.replace("“", '"').replace("”", '"')
           .replace("„", '"').replace("‟", '"')
           .replace("‘", "'").replace("’", "'"))

    match = re.search(r"\{.*\}", s, re.DOTALL)
    if not match:
        raise ValueError("No JSON object found in response")
    s = match.group()

    # Strip control chars (keep tab, lf, cr — they'll be escaped below).
    s = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", s)
    # Escape lone backslashes that aren't part of a valid JSON escape.
    s = re.sub(r'\\(?![\\"/bfnrtu])', r'\\\\', s)
    # Escape raw newlines/tabs that sit inside string literals.
    s = re.sub(
        r'"((?:[^"\\]|\\.)*)"',
        lambda m: '"' + m.group(1).replace('\n', '\\n').replace('\r', '\\r').replace('\t', '\\t') + '"',
        s,
    )
    # Remove trailing commas before } or ].
    s = re.sub(r",(\s*[}\]])", r"\1", s)

    # String-aware balanced trim — earlier version ignored quoted braces and
    # could close on a `}` that actually sat inside a string literal.
    depth = 0
    last_close = -1
    in_string = False
    escape = False
    for i, ch in enumerate(s):
        if escape:
            escape = False
            continue
        if ch == '\\':
            escape = True
            continue
        if ch == '"':
            in_string = not in_string
            continue
        if in_string:
            continue
        if ch == '{':
            depth += 1
        elif ch == '}':
            depth -= 1
            if depth == 0:
                last_close = i
                break
    if last_close != -1:
        s = s[: last_close + 1]
    return s

def _robust_parse(raw: str, label: str, repair_options: dict) -> dict:
    """Three-tier JSON recovery:
       1. local cleanup (handles ~95% of LLM JSON quirks),
       2. ask the model to repair generically,
       3. ask again with the specific parser error fed back in.
    """
    # Tier 1
    try:
        return json.loads(_extract_json_object(raw))
    except (ValueError, json.JSONDecodeError) as e1:
        print(f"{label}: local parse failed ({e1}), asking model to repair…")
        err = e1
        last_raw = raw

    # Tier 2 — model repair with format=json so output is grammar-constrained.
    repair = ("The following is supposed to be a JSON object but contains "
              "syntax errors. Return ONLY the corrected JSON object:\n\n"
              + last_raw[:8000])
    raw2 = ask_ollama("Output only valid JSON.", repair,
                      options={**repair_options, "temperature": 0.05},
                      force_json=True)
    try:
        return json.loads(_extract_json_object(raw2))
    except (ValueError, json.JSONDecodeError) as e2:
        print(f"{label}: tier-2 repair failed ({e2}), asking with error feedback…")
        err = e2
        last_raw = raw2

    # Tier 3 — feed the parser's error message back, again with format=json.
    repair2 = (f"Python's json.loads raised: {err}\n\n"
               "Fix that specific issue and return ONLY the corrected JSON "
               "object:\n\n" + last_raw[:8000])
    raw3 = ask_ollama("You are a JSON repair tool. Output only valid JSON.",
                      repair2,
                      options={**repair_options, "temperature": 0.0},
                      force_json=True)
    try:
        return json.loads(_extract_json_object(raw3))
    except (ValueError, json.JSONDecodeError) as e3:
        raise ValueError(f"{label}: all repair tiers failed; last error: {e3}") from e3

def _parse_json_stage(raw: str, label: str) -> dict:
    """Back-compat shim; new callers should use _robust_parse directly."""
    try:
        return json.loads(_extract_json_object(raw))
    except (ValueError, json.JSONDecodeError) as e:
        raise ValueError(f"{label}: {e}") from e

def run_research_stage(req: "WorkRequest") -> dict:
    system_prompt = _load_prompt("research_prompt.txt")

    # Extract text + page images from uploaded documents (pdf/docx/pptx/xlsx)
    doc_text, doc_images = process_documents(req.documents) if req.documents else ("", [])
    combined_material = "\n\n".join(p for p in [req.material, doc_text] if p)
    combined_images   = list(req.images) + doc_images
    print(f"Stage 1 inputs: material={len(combined_material)} chars, images={len(combined_images)}")

    image_note = ""
    if combined_images:
        image_note = (
            f"\n\nIMAGES: {len(combined_images)} image(s) attached as study material "
            "(may include rendered document pages). "
            "Read every visible text, diagram label, formula, and annotation. "
            "Treat the image content as authoritative source material — use facts from the images "
            "in core_concepts, key_facts, and misconceptions."
        )
    user_message = (
        f"Topic: {req.lesson.topic}\n"
        f"School type: {req.lesson.school_type or 'not specified'}\n"
        f"Grade: {req.lesson.grade or 'not specified'}\n"
        f"Investigation style: {req.lesson.mode}\n"
        f"Language: {req.lesson.language}\n"
        f"Material:\n{combined_material[:8000] if combined_material else 'none'}"
        f"{image_note}"
    )
    raw = ask_ollama(system_prompt, user_message,
                     options=_OPTS_RESEARCH, images=combined_images or None,
                     force_json=True)
    print(f"Stage 1 raw length: {len(raw)}")
    return _robust_parse(raw, "Stage 1 Research", _OPTS_RESEARCH)

def run_task_design_stage(req: "WorkRequest", research: dict) -> dict:
    system_prompt = _load_prompt("task_design_prompt.txt")
    user_message = json.dumps({
        "research": research,
        "lesson":   req.lesson.model_dump(),
    }, ensure_ascii=False)
    raw = ask_ollama(system_prompt, user_message,
                     options=_OPTS_TASK_DESIGN, force_json=True)
    print(f"Stage 2 raw length: {len(raw)}")
    return _robust_parse(raw, "Stage 2 Tasks", _OPTS_TASK_DESIGN)

def run_story_stage(req: "WorkRequest", research: dict, tasks: dict) -> "SessionBlueprint":
    system_prompt = _load_prompt("setup_prompt.txt")
    user_message = json.dumps({
        "teacher":        req.teacher.model_dump(),
        "lesson":         req.lesson.model_dump(),
        "material":       req.material,
        "topic_research": research,
        "task_list":      tasks,
    }, ensure_ascii=False)
    raw = ask_ollama(system_prompt, user_message,
                     options=_OPTS_BLUEPRINT, force_json=True)
    print(f"Stage 3 raw length: {len(raw)}")
    data = _robust_parse(raw, "Stage 3 Blueprint", _OPTS_BLUEPRINT)
    return SessionBlueprint._from_dict(data)

def _run_validation(blueprint: "SessionBlueprint") -> None:
    try:
        entries = "\n\n".join(
            f"[{i}] Subject: {blueprint.topic.core_subject}\n"
            f"Expected answer: {m.expected_student_response}"
            for i, m in enumerate(blueprint.moves)
        )
        val_result = ask_ollama(_VALIDATION_SYSTEM, entries, options=_OPTS_VALIDATE)
        flags = [l for l in val_result.splitlines() if l.strip().startswith("FLAG:")]
        if flags:
            print(f"Validation flags ({len(flags)}):")
            for fl in flags:
                print(" ", fl)
    except Exception as e:
        print(f"Validation skipped: {e}")

# ── Session store ────────────────────────────────────────────────

_sessions:      dict[str, SessionBlueprint] = {}
_session_state: dict[str, dict]             = {}

def _init_state(blueprint: SessionBlueprint) -> dict:
    return {
        "current_move_id":       blueprint.opening.opening_move_id,
        "visited_move_ids":      [],
        "turn_count":            0,
        "correct_count":         0,
        "closed":                False,
    }

def _phase_index_for_move(blueprint: SessionBlueprint, move_id: str) -> int:
    move = blueprint.get_move(move_id)
    if not move:
        return 0
    for i, phase in enumerate(blueprint.phases):
        if phase.id == move.phase:
            return i
    return 0

def _closing_line(blueprint: SessionBlueprint, state: dict) -> str:
    total = max(state["turn_count"], 1)
    ratio = state["correct_count"] / total
    if ratio >= 0.70:
        return blueprint.closing.success_line
    elif ratio >= 0.40:
        return blueprint.closing.partial_line
    else:
        return blueprint.closing.fail_line

def _apply_state_signal(
    signal:       str,
    current_move: "Move | None",
    blueprint:    SessionBlueprint,
    state:        dict,
) -> None:
    # Extracted from the inline handler because the original had CLOSE unreachable
    # (it was inside the ADVANCE branch). Now all three signals are first-class.
    if signal == "ADVANCE":
        state["correct_count"] += 1
        if current_move and current_move.type == "closing":
            state["closed"] = True
        elif current_move and current_move.if_correct:
            next_id = current_move.if_correct
            state["current_move_id"] = next_id
            if next_id not in state["visited_move_ids"]:
                state["visited_move_ids"].append(next_id)
    elif signal == "CLOSE":
        state["closed"] = True
    else:  # STAY
        if current_move and current_move.if_incorrect:
            state["current_move_id"] = current_move.if_incorrect

    if state["turn_count"] >= blueprint.session.estimated_turns:
        state["closed"] = True

# ── API ──────────────────────────────────────────────────────────

@app.post("/api/work")
def work(req: WorkRequest):
    print("Request received (sync)")
    try:
        research  = run_research_stage(req)
        tasks     = run_task_design_stage(req, research)
        blueprint = run_story_stage(req, research, tasks)
        _run_validation(blueprint)
        _sessions[blueprint.session_id] = blueprint
        _session_state[blueprint.session_id] = _init_state(blueprint)
        print(f"Blueprint generated: {blueprint.session_id}")
        return blueprint.to_dict()
    except ValueError as e:
        raise HTTPException(status_code=422, detail=f"Model returned invalid JSON: {e}")
    except requests.RequestException as e:
        raise HTTPException(status_code=503, detail=f"Ollama unreachable: {e}")

@app.post("/api/work/stream")
def work_stream(req: WorkRequest):
    def gen():
        try:
            yield f"data: {json.dumps({'type': 'progress', 'stage': 1, 'label': 'Analyzing topic and concepts…'})}\n\n"
            research = run_research_stage(req)
            yield f"data: {json.dumps({'type': 'stage_done', 'stage': 1})}\n\n"

            yield f"data: {json.dumps({'type': 'progress', 'stage': 2, 'label': 'Designing learning tasks…'})}\n\n"
            tasks = run_task_design_stage(req, research)
            yield f"data: {json.dumps({'type': 'stage_done', 'stage': 2})}\n\n"

            yield f"data: {json.dumps({'type': 'progress', 'stage': 3, 'label': 'Building your noir case…'})}\n\n"
            blueprint = run_story_stage(req, research, tasks)
            _run_validation(blueprint)

            _sessions[blueprint.session_id] = blueprint
            _session_state[blueprint.session_id] = _init_state(blueprint)

            yield f"data: {json.dumps({'type': 'stage_done', 'stage': 3})}\n\n"
            yield f"data: {json.dumps({'type': 'complete', 'blueprint': blueprint.to_dict()})}\n\n"

        except Exception as e:
            import traceback; traceback.print_exc()
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)})}\n\n"

    return StreamingResponse(
        gen(), media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )

def _restore_blueprint(data: dict) -> SessionBlueprint:
    """Rebuild a SessionBlueprint from a dict produced by `SessionBlueprint.to_dict()`.
    Used by /api/session/restore so the frontend can recover after a backend restart
    without having to regenerate the whole case.
    """
    def fit(cls_, d):
        if not isinstance(d, dict):
            return cls_()
        keys = {f.name for f in dataclasses.fields(cls_)}
        return cls_(**{k: v for k, v in d.items() if k in keys})
    return SessionBlueprint(
        session_id=data.get("session_id") or "session",
        teacher=fit(Teacher, data.get("teacher") or {}),
        session=fit(Session, data.get("session") or {}),
        topic=fit(Topic, data.get("topic") or {}),
        phases=[fit(Phase, p) for p in (data.get("phases") or [])],
        moves=[fit(Move, m) for m in (data.get("moves") or [])],
        characters=[fit(Character, c) for c in (data.get("characters") or [])],
        contingencies=fit(Contingencies, data.get("contingencies") or {}),
        opening=fit(Opening, data.get("opening") or {}),
        closing=fit(Closing, data.get("closing") or {}),
        execution_rules=fit(ExecutionRules, data.get("execution_rules") or {}),
    )

@app.post("/api/session/restore")
def restore_session(blueprint_data: dict):
    try:
        bp = _restore_blueprint(blueprint_data)
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Invalid blueprint: {e}")
    already = bp.session_id in _sessions
    if not already:
        _sessions[bp.session_id] = bp
        _session_state[bp.session_id] = _init_state(bp)
    return {"session_id": bp.session_id, "restored": not already}

@app.post("/api/chat")
def chat(req: ChatRequest):
    blueprint = _sessions.get(req.session_id)
    if not blueprint:
        raise HTTPException(status_code=404, detail="Session not found")

    state = _session_state.setdefault(req.session_id, _init_state(blueprint))

    if state["closed"]:
        return {
            "reply":                  _closing_line(blueprint, state),
            "phase_index":            len(blueprint.phases) - 1,
            "closed":                 True,
        }

    current_move = blueprint.get_move(state["current_move_id"])
    is_hint = req.message.startswith("[SYSTEM:")
    accuracy_ratio = state.get("correct_count", 0) / max(state.get("turn_count", 1), 1)

    try:
        reply, signal = run_teacher_agent(
            blueprint, req.message, current_move, req.history,
            visited_move_ids=state.get("visited_move_ids", []),
            accuracy_ratio=accuracy_ratio,
        )
    except requests.RequestException as e:
        raise HTTPException(status_code=503, detail=f"Ollama unreachable: {e}")

    if not is_hint:
        state["turn_count"] += 1
        if current_move and current_move.id not in state["visited_move_ids"]:
            state["visited_move_ids"].append(current_move.id)
        _apply_state_signal(signal, current_move, blueprint, state)
    phase_index = _phase_index_for_move(blueprint, state["current_move_id"])

    response: dict = {
        "reply":                  reply,
        "phase_index":            phase_index,
        "mood":                   _mood_for(blueprint.session.mode, phase_index),
        "closed":                 state["closed"],
    }
    if state["closed"]:
        response["closing_line"] = _closing_line(blueprint, state)
    return response

@app.post("/api/chat/stream")
def chat_stream(req: ChatRequest):
    blueprint = _sessions.get(req.session_id)
    if not blueprint:
        raise HTTPException(status_code=404, detail="Session not found")

    state = _session_state.setdefault(req.session_id, _init_state(blueprint))

    if state["closed"]:
        closing = _closing_line(blueprint, state)
        def closed_gen():
            yield f"data: {json.dumps({'chunk': closing})}\n\n"
            yield f"data: {json.dumps({'done': True, 'closed': True, 'closing_line': closing, 'phase_index': len(blueprint.phases) - 1})}\n\n"
        return StreamingResponse(
            closed_gen(), media_type="text/event-stream",
            headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
        )

    current_move = blueprint.get_move(state["current_move_id"])
    is_hint = req.message.startswith("[SYSTEM:")
    accuracy_ratio = state.get("correct_count", 0) / max(state.get("turn_count", 1), 1)

    system_prompt, user_message = _build_teacher_prompts(
        blueprint, req.message, current_move, req.history,
        visited_move_ids=state.get("visited_move_ids", []),
        accuracy_ratio=accuracy_ratio,
    )

    def event_gen():
        buffer = ""
        try:
            for chunk in ask_ollama_stream(system_prompt, user_message, options=_OPTS_CHAT):
                buffer += chunk
                yield f"data: {json.dumps({'chunk': chunk})}\n\n"
        except requests.RequestException as e:
            yield f"data: {json.dumps({'error': str(e)})}\n\n"
            return

        # Parse state from completed buffer
        signal = "STAY"
        state_match = re.search(
            r"\[STATE\][:\s]*\n?\s*(ADVANCE|STAY|CLOSE)", buffer, re.IGNORECASE
        )
        if state_match:
            signal = state_match.group(1).upper()

        if not is_hint:
            state["turn_count"] += 1
            if current_move and current_move.id not in state["visited_move_ids"]:
                state["visited_move_ids"].append(current_move.id)
            _apply_state_signal(signal, current_move, blueprint, state)
        phase_index = _phase_index_for_move(blueprint, state["current_move_id"])

        done_data: dict = {
            "done":                   True,
            "state_signal":           signal,
            "phase_index":            phase_index,
            "mood":                   _mood_for(blueprint.session.mode, phase_index),
            "closed":                 state["closed"],
            "accuracy_ratio":         state.get("correct_count", 0) / max(state.get("turn_count", 1), 1),
        }
        if state["closed"]:
            done_data["closing_line"] = _closing_line(blueprint, state)
        yield f"data: {json.dumps(done_data)}\n\n"

    return StreamingResponse(
        event_gen(), media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )

# ── Interrogation ────────────────────────────────────────────────

def _build_interrogation_prompt(
    blueprint: "SessionBlueprint", character: "Character",
) -> str:
    template = _load_prompt("interrogation_system_prompt.txt")
    return template.format(
        character_name=character.name or "Unknown",
        character_role=character.role or "—",
        character_personality=character.personality or "—",
        character_voice=character.voice or "—",
        character_knowledge=character.knowledge or "(no specific knowledge recorded)",
        character_secrets=character.secrets or "(none)",
        character_ignorant_of=character.ignorant_of or "(none)",
        language=blueprint.session.language or "English",
    )

def _format_interrogation_history(history: list["ChatMessage"]) -> str:
    if not history:
        return ""
    parts: list[str] = []
    for msg in history:
        label = "DETECTIVE" if msg.role == "user" else "YOU"
        parts.append(f"{label}: {msg.content}")
    return "\n".join(parts)

@app.get("/api/session/{session_id}/characters")
def list_characters(session_id: str):
    blueprint = _sessions.get(session_id)
    if not blueprint:
        raise HTTPException(status_code=404, detail="Session not found")
    return {
        "characters": [
            {
                "id":          c.id or c.name,
                "name":        c.name,
                "role":        c.role,
                "personality": c.personality,
            }
            for c in blueprint.characters
            if (c.name or c.id)
        ]
    }

@app.post("/api/interrogate")
def interrogate(req: InterrogationRequest):
    blueprint = _sessions.get(req.session_id)
    if not blueprint:
        raise HTTPException(status_code=404, detail="Session not found")

    # Match by id first, then by name (case-insensitive) for resilience.
    character = blueprint.get_character(req.character_id)
    if not character:
        for c in blueprint.characters:
            if (c.name or "").strip().lower() == req.character_id.strip().lower():
                character = c
                break
    if not character:
        raise HTTPException(status_code=404, detail="Character not found in session")

    system_prompt = _build_interrogation_prompt(blueprint, character)
    history_text  = _format_interrogation_history(req.history)
    user_message  = (
        (f"PRIOR EXCHANGE:\n{history_text}\n\n" if history_text else "")
        + f"DETECTIVE: {req.message}"
    )

    try:
        reply = ask_ollama(system_prompt, user_message, options=_OPTS_INTERROGATE)
    except requests.RequestException as e:
        raise HTTPException(status_code=503, detail=f"Ollama unreachable: {e}")

    # Strip stray section markers (but keep the content that follows them).
    reply = re.sub(r"\[(?:NARRATION|NPC:[^\]]+|QUESTION|STATE)\]\s*", "", reply).strip()
    return {"reply": reply, "character": {"id": character.id, "name": character.name}}

# ── Serve frontend ───────────────────────────────────────────────
app.mount("/", StaticFiles(directory=".", html=True), name="static")
